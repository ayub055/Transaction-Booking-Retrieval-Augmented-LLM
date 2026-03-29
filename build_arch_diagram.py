"""
Draw the Transaction Tagger technical architecture diagram as a PNG,
then insert it as a new slide in Transaction_Tagger_Presentation.pptx.

Layout (top-down):
  [INPUT]
      ↓
  [FUSION ENCODER]  (3 sub-inputs: IndicBERT | Categorical | Numeric)
      ↓
  [STRATIFIED FAISS INDEX]  ←→  [PER-CLASS PROTOTYPES]
      ↓                              ↓
  [BLENDED SCORING  α·kNN + (1-α)·Prototype]
      ↓
  [OOD DETECTION + CONFIDENCE CHECK]
      ↓              ↓
  [DIRECT OUTPUT]  [LLM FALLBACK]
      ↓              ↓
       [FINAL OUTPUT + AUDIT TRAIL]
"""

from PIL import Image, ImageDraw, ImageFont
import os, math
from pptx import Presentation
from pptx.util import Inches, Emu

# ── Canvas ────────────────────────────────────────────────────────────────────
SCALE   = 2          # 2× for retina sharpness
W_PX    = 2400 * SCALE
H_PX    = 1350 * SCALE
BG      = (247, 248, 250)

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY        = (27,  58, 107)
MID_BLUE    = (39,  91, 173)
TEAL        = (18, 120, 130)
GREEN       = (26, 122,  74)
CORAL       = (192, 57,  43)
GOLD        = (184, 134,  11)
GREY_BORDER = (204, 211, 221)
GREY_TEXT   = (85,  101, 120)
DARK_TEXT   = (26,  31,  46)
WHITE       = (255, 255, 255)
LIGHT_NAVY  = (232, 237, 247)
LIGHT_GREEN = (232, 245, 238)
LIGHT_GOLD  = (255, 248, 225)
ARROW_CLR   = (130, 145, 165)

# ── Fonts ─────────────────────────────────────────────────────────────────────
FONT_PATH  = "/System/Library/Fonts/Supplemental/Arial.ttf"
FONT_BOLD  = "/System/Library/Fonts/Helvetica.ttc"  # ttc index 0 = regular, 1 = bold

def load(path, size, index=0):
    try:
        return ImageFont.truetype(path, size * SCALE, index=index)
    except Exception:
        return ImageFont.load_default()

F_TITLE   = load(FONT_BOLD,  11, index=1)   # Helvetica Bold
F_LABEL   = load(FONT_BOLD,   9, index=1)
F_SUBLABEL= load(FONT_PATH,   8)
F_TINY    = load(FONT_PATH,   7)
F_HEAD    = load(FONT_BOLD,  13, index=1)

# ── Draw helpers ──────────────────────────────────────────────────────────────

def s(v):
    """Scale value."""
    return int(v * SCALE)


def rounded_rect(draw, x, y, w, h, r, fill, outline=None, lw=1):
    x, y, w, h, r, lw = int(x), int(y), int(w), int(h), int(r), int(lw)
    draw.rectangle([x+r, y, x+w-r, y+h], fill=fill)
    draw.rectangle([x, y+r, x+w, y+h-r], fill=fill)
    draw.ellipse([x,     y,     x+2*r, y+2*r], fill=fill)
    draw.ellipse([x+w-2*r, y,   x+w,   y+2*r], fill=fill)
    draw.ellipse([x,     y+h-2*r, x+2*r, y+h], fill=fill)
    draw.ellipse([x+w-2*r, y+h-2*r, x+w, y+h], fill=fill)
    if outline:
        draw.arc([x, y, x+2*r, y+2*r], 180, 270, fill=outline, width=lw)
        draw.arc([x+w-2*r, y, x+w, y+2*r], 270, 360, fill=outline, width=lw)
        draw.arc([x, y+h-2*r, x+2*r, y+h], 90, 180, fill=outline, width=lw)
        draw.arc([x+w-2*r, y+h-2*r, x+w, y+h], 0, 90, fill=outline, width=lw)
        draw.line([x+r, y, x+w-r, y], fill=outline, width=lw)
        draw.line([x+r, y+h, x+w-r, y+h], fill=outline, width=lw)
        draw.line([x, y+r, x, y+h-r], fill=outline, width=lw)
        draw.line([x+w, y+r, x+w, y+h-r], fill=outline, width=lw)


def centered_text(draw, text, cx, cy, font, color=DARK_TEXT):
    bbox = draw.textbbox((0, 0), text, font=font)
    tw   = bbox[2] - bbox[0]
    th   = bbox[3] - bbox[1]
    draw.text((cx - tw//2, cy - th//2), text, font=font, fill=color)


def wrapped_text(draw, text, cx, cy, max_w, font, color=DARK_TEXT, line_gap=None):
    """Center-aligned word-wrapped text."""
    words = text.split()
    lines = []
    cur   = []
    for w in words:
        test = " ".join(cur + [w])
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] > max_w and cur:
            lines.append(" ".join(cur))
            cur = [w]
        else:
            cur.append(w)
    if cur:
        lines.append(" ".join(cur))

    lh = (draw.textbbox((0, 0), "Ag", font=font)[3] + (line_gap or s(2)))
    total_h = lh * len(lines)
    y = cy - total_h // 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=font)
        tw   = bbox[2] - bbox[0]
        draw.text((cx - tw//2, y), line, font=font, fill=color)
        y += lh


def arrow(draw, x1, y1, x2, y2, color=ARROW_CLR, width=2, head=8):
    width, head = int(width * SCALE), int(head * SCALE)
    draw.line([int(x1), int(y1), int(x2), int(y2)], fill=color, width=width)
    # arrowhead
    angle = math.atan2(y2 - y1, x2 - x1)
    for da in (0.45, -0.45):
        ex = x2 - head * math.cos(angle - da)
        ey = y2 - head * math.sin(angle - da)
        draw.line([int(x2), int(y2), int(ex), int(ey)], fill=color, width=width)


def horiz_connector(draw, x1, y_mid, x2, color=ARROW_CLR, width=1.5):
    """Draw a horizontal dashed line (no arrowhead)."""
    draw.line([int(x1), int(y_mid), int(x2), int(y_mid)],
              fill=color, width=int(width * SCALE))


# ── Box factory ───────────────────────────────────────────────────────────────

def box(draw, cx, cy, w, h,
        title, subtitle=None, badge=None,
        fill=WHITE, outline=GREY_BORDER, title_color=DARK_TEXT,
        sub_color=GREY_TEXT, r=10, lw=2,
        header_fill=None, header_h=None):
    """
    Draw a rounded box. If header_fill given, top strip uses header_fill.
    Returns (cx, top_y, bottom_y, left_x, right_x).
    """
    x, y = cx - w//2, cy - h//2
    r_px = s(r)
    lw_px= int(lw * SCALE)
    rounded_rect(draw, x, y, w, h, r_px, fill, outline, lw_px)

    if header_fill and header_h:
        hh = int(header_h)
        rounded_rect(draw, x, y, w, hh, r_px, header_fill)
        # fix bottom corners of header to be square
        draw.rectangle([x, y + hh - r_px, x + w, y + hh], fill=header_fill)
        # left accent bar
        draw.rectangle([x, y, x + s(5), y + hh], fill=darken(header_fill, 30))
        # header title
        wrapped_text(draw, title, cx, y + hh//2, w - s(20), F_LABEL,
                     color=WHITE)
        # subtitle below header
        if subtitle:
            sub_y = y + hh + (h - hh)//2
            wrapped_text(draw, subtitle, cx, sub_y, w - s(20), F_SUBLABEL,
                         color=sub_color)
    else:
        # just centered text
        text_cy = cy if not subtitle else cy - s(7)
        wrapped_text(draw, title, cx, text_cy, w - s(16), F_LABEL,
                     color=title_color)
        if subtitle:
            wrapped_text(draw, subtitle, cx, cy + s(9), w - s(16), F_TINY,
                         color=sub_color)

    if badge:
        bx, by = x + w - s(36), y - s(10)
        rounded_rect(draw, bx, by, s(34), s(18), s(9), GOLD)
        centered_text(draw, badge, bx + s(17), by + s(9), F_TINY, WHITE)

    return cx, y, y + h, x, x + w


def darken(color, amount):
    return tuple(max(0, c - amount) for c in color)


# ── Build the diagram ─────────────────────────────────────────────────────────

def build():
    img  = Image.new("RGB", (W_PX, H_PX), BG)
    draw = ImageDraw.Draw(img)

    # ── Slide header ─────────────────────────────────────────────────────────
    draw.rectangle([0, 0, W_PX, s(4)], fill=NAVY)
    draw.rectangle([0, 0, s(7), H_PX], fill=NAVY)
    centered_text(draw, "TECHNICAL ARCHITECTURE",
                  W_PX//2, s(28), F_HEAD, NAVY)
    # thin gold underline
    draw.rectangle([s(30), s(45), W_PX - s(30), s(47)], fill=GREY_BORDER)

    # ────────────────────────────────────────────────────────────────────────
    # Coordinate helpers (all in pixel coords, already scaled via s())
    # Layout column centres (left → right: 5 columns)
    # Column layout (horizontal pipeline):
    #   C0: Input
    #   C1: Fusion Encoder (with 3 sub-nodes above)
    #   C2: Retrieval layer (2 sub-nodes: FAISS + Prototypes)
    #   C3: Scoring + Decision
    #   C4: Output
    # ────────────────────────────────────────────────────────────────────────

    PAD_L  = s(55)
    PAD_R  = s(55)
    USE_W  = W_PX - PAD_L - PAD_R   # usable width
    PAD_T  = s(58)
    PAD_B  = s(30)
    USE_H  = H_PX - PAD_T - PAD_B

    N_COL  = 5
    col_w  = USE_W // N_COL
    col_cx = [PAD_L + col_w * i + col_w // 2 for i in range(N_COL)]

    # Row vertical centres (4 rows)
    ROW_H  = USE_H // 4
    row_cy = [PAD_T + ROW_H * i + ROW_H // 2 for i in range(4)]

    BW  = int(col_w * 0.78)   # box width
    BH1 = s(60)               # main box height
    BH2 = s(50)               # sub box height

    # ── ROW 0: Input fields (3 boxes) ────────────────────────────────────────
    # Centred in cols 0..2
    inp_fields = [
        ("Transaction\nDescription", "Free-text, mixed language\n(Hindi/English/UPI IDs)"),
        ("Payment\nMode", "UPI · NEFT · IMPS · RTGS\n(Categorical)"),
        ("Amount\n(₹)", "Raw value →\nStandardScaler"),
    ]
    inp_cx = [col_cx[0], col_cx[1], col_cx[2]]
    inp_boxes = []
    for (title, sub), cx in zip(inp_fields, inp_cx):
        b = box(draw, cx, row_cy[0], BW, BH1,
                title, sub,
                fill=LIGHT_NAVY, outline=MID_BLUE, title_color=NAVY,
                sub_color=GREY_TEXT, r=8, lw=2,
                header_fill=MID_BLUE, header_h=s(24))
        inp_boxes.append(b)

    # ── Encoder box (row 1, col 0..2, spans 3 columns) ───────────────────────
    enc_cx  = (col_cx[0] + col_cx[2]) // 2
    enc_w   = col_cx[2] + BW//2 - (col_cx[0] - BW//2)
    enc_y   = row_cy[1]
    enc_box = box(draw, enc_cx, enc_y, enc_w, BH1,
                  "FUSION ENCODER",
                  "IndicBERT text (768-dim)  +  Categorical embedding (4-dim)  +  Numeric  →  MLP + LayerNorm  →  256-dim L2-normalised vector",
                  fill=WHITE, outline=NAVY, title_color=NAVY,
                  sub_color=GREY_TEXT, r=8, lw=2,
                  header_fill=NAVY, header_h=s(24))

    # arrows: input fields → encoder
    for b in inp_boxes:
        arrow(draw, b[0], b[1] + BH1,    # bottom of input box
              b[0], enc_box[1],           # top of encoder box
              color=MID_BLUE, width=1.5)

    # ── FAISS box (row 2, col 0..1) ───────────────────────────────────────────
    faiss_cx = (col_cx[0] + col_cx[1]) // 2
    faiss_w  = col_cx[1] + BW//2 - (col_cx[0] - BW//2)
    faiss_y  = row_cy[2]
    faiss_box= box(draw, faiss_cx, faiss_y, faiss_w, BH1,
                   "STRATIFIED FAISS INDEX",
                   "HNSW graph  ·  Max 500 exemplars/class\n(greedy farthest-point sampling)  ·  Top-10 k-NN retrieval",
                   fill=LIGHT_NAVY, outline=MID_BLUE, title_color=NAVY,
                   sub_color=GREY_TEXT, r=8, lw=2,
                   header_fill=MID_BLUE, header_h=s(24))

    # ── Prototype box (row 2, col 2) ──────────────────────────────────────────
    proto_cx = col_cx[2]
    proto_box= box(draw, proto_cx, row_cy[2], BW, BH1,
                   "PER-CLASS PROTOTYPES",
                   "L2-normalised centroids\n(one per category, 43 total)",
                   fill=LIGHT_NAVY, outline=MID_BLUE, title_color=NAVY,
                   sub_color=GREY_TEXT, r=8, lw=2,
                   header_fill=MID_BLUE, header_h=s(24))

    # arrow encoder → faiss + proto (split)
    enc_bot_y = enc_box[1] + BH1
    split_x   = enc_cx
    split_y   = enc_bot_y + s(14)
    draw.line([split_x, enc_bot_y, split_x, split_y],
              fill=ARROW_CLR, width=int(1.5 * SCALE))
    # horizontal to faiss
    arrow(draw, split_x, split_y, faiss_cx, faiss_box[1],
          color=MID_BLUE, width=1.5)
    # horizontal to proto
    arrow(draw, split_x, split_y, proto_cx, proto_box[1],
          color=MID_BLUE, width=1.5)

    # ── Blended scoring box (row 2, col 3) ────────────────────────────────────
    score_cx = col_cx[3]
    score_y  = row_cy[2]
    score_box= box(draw, score_cx, score_y, BW, BH1,
                   "BLENDED SCORING",
                   "α · kNN-vote  +  (1-α) · prototype\nα=0.6  ·  distance-weighted",
                   fill=LIGHT_GREEN, outline=GREEN, title_color=GREEN,
                   sub_color=GREY_TEXT, r=8, lw=2,
                   header_fill=GREEN, header_h=s(24))

    # arrows faiss → scoring, proto → scoring
    arrow(draw, faiss_cx + faiss_w//2, faiss_y,
          score_cx - BW//2, score_y,
          color=GREEN, width=1.5)
    arrow(draw, proto_cx + BW//2, row_cy[2],
          score_cx - BW//2, score_y,
          color=GREEN, width=1.5)

    # ── OOD + Confidence box (row 2, col 4) ───────────────────────────────────
    ood_cx  = col_cx[4]
    ood_box = box(draw, ood_cx, row_cy[2], BW, BH1,
                  "OOD DETECTION",
                  "Distance threshold (95th pctl)\nCalibrated on real 'other' txns\n< 0.75 conf → escalate",
                  fill=LIGHT_GOLD, outline=GOLD, title_color=GOLD,
                  sub_color=GREY_TEXT, r=8, lw=2,
                  header_fill=GOLD, header_h=s(24))

    # arrow scoring → OOD
    arrow(draw, score_cx + BW//2, score_y,
          ood_cx - BW//2, row_cy[2],
          color=GOLD, width=1.5)

    # ── OUTPUT row (row 3) ────────────────────────────────────────────────────
    # Left output (direct)
    out_cx  = col_cx[1]
    out_box = box(draw, out_cx, row_cy[3], BW, BH1,
                  "DIRECT OUTPUT",
                  "Predicted category\nConfidence score\nTop-3 audit trail",
                  fill=LIGHT_GREEN, outline=GREEN, title_color=GREEN,
                  sub_color=GREY_TEXT, r=8, lw=2,
                  header_fill=GREEN, header_h=s(24),
                  badge="≥0.75 conf")

    # Right output (LLM fallback)
    llm_cx  = col_cx[3]
    llm_box = box(draw, llm_cx, row_cy[3], BW, BH1,
                  "LLM FALLBACK",
                  "Ollama / OpenAI / Anthropic\nFew-shot prompt with top-3 neighbours\n~10% of transactions",
                  fill=LIGHT_GOLD, outline=GOLD, title_color=GOLD,
                  sub_color=GREY_TEXT, r=8, lw=2,
                  header_fill=GOLD, header_h=s(24),
                  badge="OOD / low conf")

    # arrows OOD → output + LLM
    ood_bot_y = row_cy[2] + BH1//2
    arrow(draw, ood_cx, ood_bot_y + s(10),
          out_cx,  row_cy[3] - BH1//2,
          color=GREEN, width=1.5)
    arrow(draw, ood_cx, ood_bot_y + s(10),
          llm_cx, row_cy[3] - BH1//2,
          color=GOLD, width=1.5)

    # ── Final output bar (bottom) ─────────────────────────────────────────────
    final_y  = row_cy[3] + BH1//2 + s(18)
    final_cx = W_PX // 2
    final_w  = int(USE_W * 0.55)
    rounded_rect(draw, final_cx - final_w//2, final_y,
                 final_w, s(34), s(8), NAVY)
    centered_text(draw,
                  "FINAL OUTPUT  —  category  ·  confidence  ·  OOD flag  ·  top-3 similar transactions",
                  final_cx, final_y + s(17), F_SUBLABEL, WHITE)

    # arrows out + llm → final
    arrow(draw, out_cx,  row_cy[3] + BH1//2,
          final_cx - final_w//4, final_y,
          color=NAVY, width=1.5)
    arrow(draw, llm_cx,  row_cy[3] + BH1//2,
          final_cx + final_w//4, final_y,
          color=NAVY, width=1.5)

    # ── Legend (top-right) ────────────────────────────────────────────────────
    leg_items = [
        (MID_BLUE, "Encoding / Retrieval layer"),
        (GREEN,    "High-confidence path (≥0.75)"),
        (GOLD,     "OOD / Low-confidence path"),
        (NAVY,     "Core infrastructure"),
    ]
    lx, ly = W_PX - s(255), s(55)
    for i, (color, label) in enumerate(leg_items):
        y0 = ly + i * s(18)
        rounded_rect(draw, lx, y0, s(14), s(12), s(3), color)
        draw.text((lx + s(19), y0 - s(1)), label, font=F_TINY, fill=GREY_TEXT)

    # ── Footer ────────────────────────────────────────────────────────────────
    draw.rectangle([0, H_PX - s(26), W_PX, H_PX], fill=LIGHT_NAVY)
    draw.rectangle([0, H_PX - s(26), W_PX, H_PX - s(25)], fill=GREY_BORDER)
    centered_text(draw,
                  "FusionEncoder (BERT+Categorical+Numeric)  ·  Stratified FAISS HNSW  ·  "
                  "SupCon Training  ·  OOD Calibration  ·  LLM Hybrid",
                  W_PX//2, H_PX - s(13), F_TINY, GREY_TEXT)

    return img


# ── Save image + inject into PPTX ─────────────────────────────────────────────

def main():
    print("Drawing architecture diagram ...")
    img = build()
    img_path = "arch_diagram.png"
    img.save(img_path, dpi=(200 * SCALE, 200 * SCALE))
    print(f"  Saved {img_path}  ({img.size[0]}×{img.size[1]} px)")

    print("Injecting into presentation ...")
    pptx_path = "Transaction_Tagger_Presentation.pptx"
    prs = Presentation(pptx_path)

    # insert after slide 2 (index 2 = slide 3 / "Our Solution")
    # We want it between slide 3 (Challenges) and slide 4 (Impact)
    # Current order: 0=Title, 1=Shift, 2=Solution, 3=Challenges, 4=Impact
    # Insert at position 3 (before Challenges)

    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    # Fill background
    bg = new_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(247, 248, 250)

    # Add image centred on slide
    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)
    img_w   = Inches(12.8)
    img_h   = Inches(img_w * H_PX / W_PX)
    img_l   = (slide_w - img_w) / 2
    img_t   = (slide_h - img_h) / 2
    new_slide.shapes.add_picture(img_path, img_l, img_t, img_w, img_h)

    # Move new slide to position 3 (after "Our Solution", before "Challenges")
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    # last slide is the one we just added (index -1)
    new_el = slides_list[-1]
    xml_slides.remove(new_el)
    xml_slides.insert(3, new_el)   # position 3 = 4th slide (0-indexed)

    prs.save(pptx_path)
    print(f"  Architecture slide inserted at position 4 in {pptx_path}")
    print("\nDone.")


if __name__ == "__main__":
    main()
