"""
Generate C-Suite presentation — Transaction Tagger project.
Design: White background, clean consulting style, human feel.
Usage:  python create_presentation.py
Output: Transaction_Tagger_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn  # noqa: E402
from lxml import etree

# ── Palette (white-background, consulting style) ──────────────────────────────
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF7, 0xF8, 0xFA)   # slide background
NAVY         = RGBColor(0x1B, 0x3A, 0x6B)   # primary brand navy
NAVY_LIGHT   = RGBColor(0xE8, 0xED, 0xF7)   # very light navy tint for fills
MID_BLUE     = RGBColor(0x27, 0x5B, 0xAD)   # table header / accents
CORAL        = RGBColor(0xC0, 0x39, 0x2B)   # old-way red
CORAL_LIGHT  = RGBColor(0xFD, 0xED, 0xEB)   # old-way row fill
GREEN        = RGBColor(0x1A, 0x7A, 0x4A)   # new-way green
GREEN_LIGHT  = RGBColor(0xE8, 0xF5, 0xEE)   # new-way row fill
GOLD         = RGBColor(0xB8, 0x86, 0x0B)   # highlight / callouts
GREY_BORDER  = RGBColor(0xCC, 0xD3, 0xDD)   # table borders
GREY_TEXT    = RGBColor(0x55, 0x65, 0x78)   # body text
DARK_TEXT    = RGBColor(0x1A, 0x1F, 0x2E)   # headings
ROW_ALT      = RGBColor(0xF2, 0xF4, 0xF8)   # alternating table row

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Core helpers ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=OFF_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line_color=None, line_w_pt=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w_pt or 0.75)
    else:
        s.line.fill.background()
    return s


def txb(slide, text, l, t, w, h,
        size=12, bold=False, italic=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return shape


def multiline_txb(slide, lines, l, t, w, h,
                  default_size=12, default_color=DARK_TEXT,
                  default_align=PP_ALIGN.LEFT, font="Calibri"):
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get("align", default_align)
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        if ln.get("space_after"):
            p.space_after = Pt(ln["space_after"])
        run = p.add_run()
        run.text = ln.get("text", "")
        run.font.size = Pt(ln.get("size", default_size))
        run.font.bold = ln.get("bold", False)
        run.font.italic = ln.get("italic", False)
        run.font.color.rgb = ln.get("color", default_color)
        run.font.name = font
    return shape


def header_band(slide, title, subtitle=None, slide_num=None):
    """Standard white-bg slide header: left navy bar + title + optional subtitle."""
    # left accent bar
    rect(slide, 0, 0, Inches(0.07), SLIDE_H, NAVY)
    # thin top line
    rect(slide, 0, 0, SLIDE_W, Inches(0.04), NAVY)

    # slide number top-right
    if slide_num:
        txb(slide, slide_num,
            SLIDE_W - Inches(0.9), Inches(0.12), Inches(0.75), Inches(0.35),
            size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)

    # title
    txb(slide, title,
        Inches(0.3), Inches(0.18), Inches(11.5), Inches(0.55),
        size=24, bold=True, color=NAVY)

    # hairline under title
    rect(slide, Inches(0.3), Inches(0.78), Inches(12.7), Pt(1.2), GREY_BORDER)

    # subtitle
    if subtitle:
        txb(slide, subtitle,
            Inches(0.3), Inches(0.84), Inches(11), Inches(0.4),
            size=12, italic=True, color=GREY_TEXT)


def footer(slide, text):
    rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), NAVY_LIGHT)
    rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Pt(1), GREY_BORDER)
    txb(slide, text,
        Inches(0.3), SLIDE_H - Inches(0.34), Inches(12.7), Inches(0.3),
        size=8, color=GREY_TEXT)


# ── Native PPTX table helper ───────────────────────────────────────────────────

def add_table(slide, rows, cols, l, t, w, h):
    """Add a native pptx table and return it."""
    tbl_shape = slide.shapes.add_table(rows, cols, l, t, w, h)
    return tbl_shape.table


def style_cell(cell, text,
               fill_color=None, text_color=DARK_TEXT,
               font_size=11, bold=False, italic=False,
               align=PP_ALIGN.LEFT, font="Calibri",
               border_color=GREY_BORDER, border_pt=0.5):
    """Apply fill, text, and border to a table cell."""
    # fill
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    else:
        cell.fill.background()

    # text
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # clear existing runs
    for run in p.runs:
        run.text = ""
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = text_color
    run.font.name = font

    # borders via XML
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ("lnL", "lnR", "lnT", "lnB"):
        ln = etree.SubElement(tcPr, qn("a:" + side))
        ln.set("w", str(int(border_pt * 12700)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{border_color[0]:02X}{border_color[1]:02X}{border_color[2]:02X}")
    return cell


def set_col_width(table, col_idx, width):
    tbl = table._tbl
    tblGrid = tbl.find(qn("a:tblGrid"))
    gridCols = tblGrid.findall(qn("a:gridCol"))
    gridCols[col_idx].set("w", str(int(width)))


def set_row_height(table, row_idx, height):
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    trs[row_idx].set("h", str(int(height)))


# ── TITLE SLIDE ───────────────────────────────────────────────────────────────

def build_title_slide(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)

    # left navy panel (full height)
    rect(slide, 0, 0, Inches(4.6), SLIDE_H, NAVY)

    # subtle diagonal accent strip inside left panel
    rect(slide, 0, Inches(5.5), Inches(4.6), Inches(0.06), MID_BLUE)
    rect(slide, 0, Inches(5.6), Inches(4.6), Inches(0.03), GOLD)

    # "AI" watermark in left panel
    txb(slide, "AI",
        Inches(0.1), Inches(1.0), Inches(4.4), Inches(4),
        size=220, bold=True,
        color=RGBColor(0x22, 0x4A, 0x85),
        align=PP_ALIGN.CENTER)

    # top-left label
    txb(slide, "CONFIDENTIAL  ·  AI STRATEGY BRIEFING",
        Inches(0.25), Inches(0.22), Inches(4.1), Inches(0.3),
        size=8, color=RGBColor(0x88, 0xA8, 0xD8))

    # right panel — title block
    multiline_txb(slide, [
        {"text": "Intelligence-Led", "size": 40, "bold": True,
         "color": NAVY},
        {"text": "Transaction", "size": 40, "bold": True,
         "color": NAVY},
        {"text": "Classification", "size": 40, "bold": True,
         "color": MID_BLUE},
    ], Inches(5.0), Inches(1.2), Inches(7.8), Inches(2.8))

    # gold rule
    rect(slide, Inches(5.0), Inches(4.05), Inches(3.0), Pt(3), GOLD)

    # subtitle
    txb(slide,
        "How the AI Team adapted Amazon Science research\n"
        "to solve a real-world Indian banking problem at scale.",
        Inches(5.0), Inches(4.2), Inches(7.8), Inches(0.9),
        size=13, italic=True, color=GREY_TEXT)

    # stat row
    stats = [
        ("43", "Categories"),
        ("35K+", "Labeled Records"),
        ("10M+", "Txns / Month"),
        ("<15ms", "Per Transaction"),
    ]
    pill_w = Inches(1.8)
    pill_h = Inches(0.82)
    gap    = Inches(0.18)
    sx     = Inches(5.0)
    sy     = Inches(5.4)
    for i, (val, lbl) in enumerate(stats):
        px = sx + i * (pill_w + gap)
        rect(slide, px, sy, pill_w, pill_h, NAVY_LIGHT, GREY_BORDER, 0.75)
        rect(slide, px, sy, Inches(0.05), pill_h, NAVY)
        multiline_txb(slide, [
            {"text": val,  "size": 20, "bold": True,  "color": NAVY,
             "align": PP_ALIGN.CENTER},
            {"text": lbl,  "size": 9,  "bold": False, "color": GREY_TEXT,
             "align": PP_ALIGN.CENTER},
        ], px, sy + Inches(0.06), pill_w, pill_h)

    # bottom-right footnote
    txb(slide,
        "Adapted from: 'Cash Transaction Booking via Retrieval-Augmented LLM'  —  Amazon Science",
        Inches(5.0), SLIDE_H - Inches(0.45), Inches(8.0), Inches(0.35),
        size=8, italic=True, color=RGBColor(0x99, 0xA8, 0xBB))

    return slide


# ── SLIDE 1 — Strategic Shift ──────────────────────────────────────────────────

def build_slide1(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    header_band(slide,
                "The Strategic Shift",
                "From brittle rule-based regex to self-learning AI",
                slide_num="01 / 04")

    # ── Two-column comparison table (native) ─────────────────────────────────
    tbl_l = Inches(0.35)
    tbl_t = Inches(1.4)
    tbl_w = Inches(12.63)
    tbl_h = Inches(5.7)
    n_data_rows = 6
    n_rows = 1 + n_data_rows   # 1 header + 6 data

    table = add_table(slide, n_rows, 2, tbl_l, tbl_t, tbl_w, tbl_h)

    # column widths (equal halves)
    half = int(tbl_w / 2)
    set_col_width(table, 0, half)
    set_col_width(table, 1, half)

    # row heights
    set_row_height(table, 0, int(Inches(0.55)))
    for r in range(1, n_rows):
        set_row_height(table, r, int(Inches(0.82)))

    # header row
    style_cell(table.cell(0, 0),
               "THE OLD WAY  —  Regex / Rule Engine",
               fill_color=CORAL, text_color=WHITE,
               font_size=13, bold=True, align=PP_ALIGN.CENTER,
               border_color=CORAL)
    style_cell(table.cell(0, 1),
               "THE NEW WAY  —  RAG + Embeddings + LLM",
               fill_color=GREEN, text_color=WHITE,
               font_size=13, bold=True, align=PP_ALIGN.CENTER,
               border_color=GREEN)

    old_rows = [
        "Hundreds of manually curated regex patterns per category",
        "Fails silently on unseen formats — new UPI apps, merchant name changes",
        "Every new payment type requires a developer to update the rules",
        "High False Positives — e.g. 'AMAZON REFUND' tagged as Shopping",
        "No concept of 'unknown' — every transaction is force-tagged into a bucket",
        "Brittle by design: one merchant rename breaks an entire category",
    ]
    new_rows = [
        "Learns from 35,000+ labeled examples; generalises to unseen formats",
        "Handles Hindi/English mix, UPI IDs, IMPS/NEFT references natively",
        "New payment types flagged as Out-of-Distribution — never silently mis-tagged",
        "Confidence-gated: uncertain predictions escalate to LLM instead of guessing",
        "Self-improving — more labelled data improves the model; zero regex changes",
        "Full audit trail: top-3 most similar past transactions per prediction",
    ]

    for r in range(n_data_rows):
        fill = CORAL_LIGHT if r % 2 == 0 else WHITE
        style_cell(table.cell(r + 1, 0),
                   f"✗   {old_rows[r]}",
                   fill_color=fill, text_color=CORAL,
                   font_size=11.5, border_color=GREY_BORDER)
        fill2 = GREEN_LIGHT if r % 2 == 0 else WHITE
        style_cell(table.cell(r + 1, 1),
                   f"✓   {new_rows[r]}",
                   fill_color=fill2, text_color=GREEN,
                   font_size=11.5, border_color=GREY_BORDER)

    footer(slide,
           "43 categories  ·  35,000+ labeled records  ·  10M+ transactions / month")
    return slide


# ── SLIDE 2 — Our Solution (3-column table) ────────────────────────────────────

def build_slide2(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    header_band(slide,
                "Our Solution: The Amazon Adaptation",
                "We didn't just apply the research — we extended it for Indian banking at scale",
                slide_num="02 / 04")

    # ── Pipeline mini-bar ────────────────────────────────────────────────────
    steps = ["Transaction Input", "Fusion Encoder", "FAISS Retrieval",
             "Vote + Prototype", "Confidence Check", "Output / LLM"]
    box_w = Inches(1.86)
    box_h = Inches(0.42)
    bx    = Inches(0.35)
    by    = Inches(1.38)
    for i, label in enumerate(steps):
        fill_c = NAVY if i < 5 else GOLD
        txt_c  = WHITE
        rect(slide, bx + i * box_w, by, box_w - Inches(0.04), box_h, fill_c)
        txb(slide, label,
            bx + i * box_w, by, box_w - Inches(0.04), box_h,
            size=9, bold=True, color=txt_c, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txb(slide, "▶",
                bx + (i + 1) * box_w - Inches(0.13), by + Inches(0.07),
                Inches(0.16), Inches(0.28),
                size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 3-column innovations table ───────────────────────────────────────────
    tbl_l = Inches(0.35)
    tbl_t = Inches(2.0)
    tbl_w = Inches(12.63)
    tbl_h = Inches(4.95)
    n_data = 5
    n_rows = 1 + n_data

    col_labels = ["INNOVATION", "GAP WE ADDRESSED", "BUSINESS IMPACT"]
    col_ws_in  = [Inches(3.1), Inches(4.77), Inches(4.76)]

    table = add_table(slide, n_rows, 3, tbl_l, tbl_t, tbl_w, tbl_h)
    for ci, cw in enumerate(col_ws_in):
        set_col_width(table, ci, int(cw))
    set_row_height(table, 0, int(Inches(0.5)))
    for r in range(1, n_rows):
        set_row_height(table, r, int(Inches(0.89)))

    # header
    for ci, lbl in enumerate(col_labels):
        style_cell(table.cell(0, ci), lbl,
                   fill_color=NAVY, text_color=WHITE,
                   font_size=11, bold=True, align=PP_ALIGN.CENTER,
                   border_color=NAVY)

    rows_data = [
        ("Stratified FAISS Index\n(class cap: 500 per category)",
         "Amazon's paper assumed balanced data.\nOur dataset is 20:1 skewed — majority classes drown minority ones.",
         "Loan Repayment, Insurance, Tax categories now achieve fair retrieval accuracy"),
        ("Prototype Blending\n(k-NN score + class centroid)",
         "Nearest-neighbour alone is noisy at decision boundaries for classes with only 250 samples.",
         "Smoother, more reliable predictions on rare transaction types — no bias toward frequent classes"),
        ("OOD Detection\n(distance-threshold calibration)",
         "The Amazon paper had no mechanism to say 'I don't know'. Every transaction was force-classified.",
         "Zero silent misclassifications. New payment types (new UPI apps, wallets) are flagged, not guessed"),
        ("IndicBERT\n(Indian language model)",
         "Amazon used English BERT trained on e-commerce text. Our data is Hindi/English transliterated.",
         "+3–5% accuracy on UPI/IMPS descriptions. Handles 'PAYTM/9876543@okaxis' natively"),
        ("LLM Fallback\n(confidence-gated escalation)",
         "Low-confidence predictions below 75% threshold need a second opinion before committing.",
         "90% of transactions are fully automated. Only ~10% escalate to LLM — keeping costs minimal"),
    ]

    for r, (c1, c2, c3) in enumerate(rows_data):
        fill = ROW_ALT if r % 2 == 0 else WHITE
        style_cell(table.cell(r + 1, 0), c1,
                   fill_color=fill, text_color=NAVY,
                   font_size=10.5, bold=True, border_color=GREY_BORDER)
        style_cell(table.cell(r + 1, 1), c2,
                   fill_color=fill, text_color=GREY_TEXT,
                   font_size=10.5, border_color=GREY_BORDER)
        style_cell(table.cell(r + 1, 2), c3,
                   fill_color=fill, text_color=GREEN,
                   font_size=10.5, border_color=GREY_BORDER)

    footer(slide,
           "15ms latency per transaction  ·  1M transactions processed in 15 minutes  ·  90% fully automated")
    return slide


# ── SLIDE 3 — Real-World Challenges ───────────────────────────────────────────

def build_slide3(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    header_band(slide,
                "Challenges Amazon's Paper Never Had to Face",
                "Research benchmarks don't have production-scale imbalance, Indian text, or budget constraints",
                slide_num="03 / 04")

    # ── 3-column challenge table: Challenge | Our Fix | Business Impact ──────
    tbl_l = Inches(0.35)
    tbl_t = Inches(1.4)
    tbl_w = Inches(12.63)
    tbl_h = Inches(5.65)
    n_challenges = 4
    n_rows = 1 + n_challenges

    col_labels = ["CHALLENGE", "HOW WE SOLVED IT", "WHY IT MATTERS TO THE BUSINESS"]
    col_ws_in  = [Inches(3.3), Inches(4.57), Inches(4.76)]

    table = add_table(slide, n_rows, 3, tbl_l, tbl_t, tbl_w, tbl_h)
    for ci, cw in enumerate(col_ws_in):
        set_col_width(table, ci, int(cw))
    set_row_height(table, 0, int(Inches(0.5)))
    for r in range(1, n_rows):
        set_row_height(table, r, int(Inches(1.29)))

    # header row
    for ci, lbl in enumerate(col_labels):
        style_cell(table.cell(0, ci), lbl,
                   fill_color=NAVY, text_color=WHITE,
                   font_size=11, bold=True, align=PP_ALIGN.CENTER,
                   border_color=NAVY)

    challenges = [
        (
            "Extreme Class Imbalance\n\n'Salary' = 5,000 samples\n'Loan Repayment' = 250 samples\n(20:1 ratio)",
            "Supervised Contrastive Loss (SupCon) replaces triplet loss — uses every in-batch positive pair.\n"
            "Combined with inverse-frequency class weighting and class-balanced batch sampling.",
            "Without this fix, 8–10 critical categories (Loans, Taxes, Insurance) would have near-zero accuracy.\n"
            "These categories carry the highest regulatory and financial risk."
        ),
        (
            "The 'Confident But Wrong' Problem\n\nNew UPI apps and wallets emerge weekly.\nNo 'unknown' mode → silent mis-tagging.",
            "Distance-based OOD detection calibrated on 192 rows of real customer transactions.\n"
            "Flags uncertainty instead of guessing — uncertain cases escalate to LLM.",
            "Eliminates the most dangerous failure mode in automated financial tagging.\n"
            "Regulators and auditors require that unknown transactions are escalated, not silently filed."
        ),
        (
            "Indian Language & Payment Stack\n\nAmazon trained on English e-commerce.\nOur data: UPI IDs, IMPS codes, Hindi merchant names in Roman script.",
            "Fine-tuned IndicBERT (trained on Indian-language text) using gradual layer unfreezing.\n"
            "Preserves general language knowledge while adapting to Indian banking domain.",
            "15–20% accuracy gap on UPI-heavy transactions without this.\n"
            "India-specific payment infrastructure is not covered by any standard English NLP model."
        ),
        (
            "Inference at 10M+ Transactions/Month on a Startup Budget\n\nNaive BERT + flat search: ~8 hours per 1M txns.\nCloud GPU at naive scale: ~$2,000/month.",
            "HNSW index: 50× faster similarity search.\n"
            "FP16 precision encoding: 2× faster, 40% less GPU memory.\n"
            "Batched GPU matrix operations for prototype scoring.",
            "1M transactions processed in 15 minutes.\n"
            "Cloud GPU cost reduced from ~$2,000/month to ~$80/month.\n"
            "Linear cost scaling — handles 10× growth without architecture changes."
        ),
    ]

    row_fills = [NAVY_LIGHT, WHITE, NAVY_LIGHT, WHITE]

    for r, (c1, c2, c3) in enumerate(challenges):
        fill = row_fills[r]
        style_cell(table.cell(r + 1, 0), c1,
                   fill_color=fill, text_color=NAVY,
                   font_size=10, bold=True, border_color=GREY_BORDER)
        style_cell(table.cell(r + 1, 1), c2,
                   fill_color=fill, text_color=GREY_TEXT,
                   font_size=10, border_color=GREY_BORDER)
        style_cell(table.cell(r + 1, 2), c3,
                   fill_color=fill, text_color=DARK_TEXT,
                   font_size=10, bold=False, border_color=GREY_BORDER)

    footer(slide,
           "SupCon loss · OOD detection · IndicBERT · HNSW index — each solving a gap the Amazon paper left open")
    return slide


# ── SLIDE 4 — Impact & Vision ──────────────────────────────────────────────────

def build_slide4(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    header_band(slide,
                "Impact & The Road Ahead",
                "From proof-of-concept to production-grade proprietary AI asset",
                slide_num="04 / 04")

    # ── Left — Before/After metrics table ────────────────────────────────────
    tbl_l = Inches(0.35)
    tbl_t = Inches(1.35)
    tbl_w = Inches(6.7)
    tbl_h = Inches(4.75)

    metric_rows = [
        ("Tagging Throughput",        "Manual / weeks",         "1M transactions in 15 min"),
        ("Unknown Transactions",      "Silent misclassification","Flagged & escalated to LLM"),
        ("Minority Class Accuracy",   "Near-zero",              "Balanced — SupCon + stratified index"),
        ("False Positive Rate",       "High (regex brittle)",   "Confidence-gated; OOD-checked"),
        ("Audit Trail",               "None",                   "Top-3 similar txns per prediction"),
        ("Cloud Infra Cost",          "N/A (manual process)",   "~$80 / month GPU"),
    ]
    n_rows = 1 + len(metric_rows)

    table = add_table(slide, n_rows, 3, tbl_l, tbl_t, tbl_w, tbl_h)
    col_ws = [Inches(2.6), Inches(2.05), Inches(2.05)]
    for ci, cw in enumerate(col_ws):
        set_col_width(table, ci, int(cw))
    set_row_height(table, 0, int(Inches(0.45)))
    for r in range(1, n_rows):
        set_row_height(table, r, int(Inches(0.72)))

    hdrs = ["METRIC", "BEFORE", "AFTER"]
    hdr_fills = [NAVY, CORAL, GREEN]
    for ci, (h, hf) in enumerate(zip(hdrs, hdr_fills)):
        style_cell(table.cell(0, ci), h,
                   fill_color=hf, text_color=WHITE,
                   font_size=10, bold=True, align=PP_ALIGN.CENTER,
                   border_color=hf)

    for r, (m, before, after) in enumerate(metric_rows):
        fill = ROW_ALT if r % 2 == 0 else WHITE
        style_cell(table.cell(r + 1, 0), m,
                   fill_color=fill, text_color=DARK_TEXT,
                   font_size=10.5, bold=True, border_color=GREY_BORDER)
        style_cell(table.cell(r + 1, 1), before,
                   fill_color=fill, text_color=CORAL,
                   font_size=10, align=PP_ALIGN.CENTER, border_color=GREY_BORDER)
        style_cell(table.cell(r + 1, 2), after,
                   fill_color=fill, text_color=GREEN,
                   font_size=10, border_color=GREY_BORDER)

    # ── Right — Roadmap ───────────────────────────────────────────────────────
    rx = Inches(7.4)
    ry = Inches(1.35)
    rw = Inches(5.6)

    txb(slide, "Strategic Roadmap",
        rx, ry, rw, Inches(0.4),
        size=13, bold=True, color=NAVY)
    rect(slide, rx, ry + Inches(0.4), Inches(2.0), Pt(2), GOLD)

    phases = [
        ("Phase 1  —  NOW",
         "43 categories · 35K labeled records · Batch inference pipeline live",
         NAVY),
        ("Phase 2  —  Active Learning",
         "Model flags uncertain cases · Analyst confirms · Dataset grows automatically",
         MID_BLUE),
        ("Phase 3  —  Real-Time Streaming",
         "Kafka / event-driven architecture for live transaction classification",
         RGBColor(0x1A, 0x7A, 0x4A)),
        ("Phase 4  —  Cross-Product Transfer",
         "Fraud detection · Credit scoring · Customer segmentation\nSame embedding infrastructure — compounding returns",
         GOLD),
    ]

    phase_h = Inches(1.06)
    for i, (title, body, color) in enumerate(phases):
        py = ry + Inches(0.55) + i * (phase_h + Inches(0.06))
        # card background
        rect(slide, rx, py, rw, phase_h, ROW_ALT if i % 2 == 0 else WHITE,
             line_color=GREY_BORDER, line_w_pt=0.75)
        # left colour accent
        rect(slide, rx, py, Inches(0.06), phase_h, color)
        # phase number circle area
        rect(slide, rx + Inches(0.14), py + Inches(0.33),
             Inches(0.32), Inches(0.32), color)
        txb(slide, str(i + 1),
            rx + Inches(0.14), py + Inches(0.33),
            Inches(0.32), Inches(0.32),
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(slide, title,
            rx + Inches(0.56), py + Inches(0.08),
            rw - Inches(0.65), Inches(0.36),
            size=11, bold=True, color=color)
        txb(slide, body,
            rx + Inches(0.56), py + Inches(0.44),
            rw - Inches(0.65), Inches(0.55),
            size=10, color=GREY_TEXT)

    # ── Closing quote ─────────────────────────────────────────────────────────
    qy = SLIDE_H - Inches(0.82)
    rect(slide, Inches(0.35), qy, Inches(12.63), Inches(0.66),
         NAVY_LIGHT, line_color=GREY_BORDER)
    rect(slide, Inches(0.35), qy, Inches(0.06), Inches(0.66), GOLD)
    txb(slide,
        '"This is not a vendor tool — it is a proprietary AI capability that '
        'compounds in value as our transaction data grows."',
        Inches(0.55), qy + Inches(0.1), Inches(12.2), Inches(0.5),
        size=11.5, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    footer(slide,
           "Adapted from: 'Cash Transaction Booking via Retrieval-Augmented LLM'  —  Amazon Science  |  "
           "Extended by the AI Team for Indian banking at scale")
    return slide


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    print("Building title slide ...")
    build_title_slide(prs)
    print("Building slide 01 — Strategic Shift ...")
    build_slide1(prs)
    print("Building slide 02 — Our Solution ...")
    build_slide2(prs)
    print("Building slide 03 — Real-World Challenges ...")
    build_slide3(prs)
    print("Building slide 04 — Impact & Vision ...")
    build_slide4(prs)

    out = "Transaction_Tagger_Presentation.pptx"
    prs.save(out)
    print(f"\nSaved -> {out}")


if __name__ == "__main__":
    main()
